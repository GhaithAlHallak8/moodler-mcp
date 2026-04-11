import json
from datetime import UTC, datetime

from moodler_mcp.client import download_file, fetch_page
from moodler_mcp.moodle_api import get_course_sections, list_enrolled_courses
from moodler_mcp.server import mcp


def _format_course(c: dict) -> dict:
    result = {
        "id": c["id"],
        "fullname": c.get("fullname", ""),
        "shortname": c.get("shortname", ""),
        "category": c.get("coursecategory", ""),
    }
    if c.get("startdate"):
        result["startdate"] = datetime.fromtimestamp(c["startdate"], tz=UTC).isoformat()
    if c.get("enddate") and c["enddate"] != 0:
        result["enddate"] = datetime.fromtimestamp(c["enddate"], tz=UTC).isoformat()
    return result


@mcp.tool()
async def list_courses(classification: str = "all", limit: int = 50) -> str:
    """List your enrolled Moodle courses.

    Args:
        classification: Filter by 'all', 'inprogress', 'past', or 'future'
        limit: Max number of courses to return (max 50)
    """
    limit = min(limit, 50)
    data = await list_enrolled_courses(
        classification=classification,
        limit=limit,
    )
    courses = [_format_course(c) for c in data.get("courses", [])]
    return json.dumps({"total": len(courses), "courses": courses}, indent=2)


@mcp.tool()
async def get_course_contents(course_id: int) -> str:
    """Get all sections, activities, and resources within a course.

    Args:
        course_id: The Moodle course ID
    """
    sections = await get_course_sections(course_id=course_id)
    return json.dumps(sections, indent=2)


_TEXT_SUFFIXES = {
    ".py",
    ".txt",
    ".csv",
    ".json",
    ".ipynb",
    ".md",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".js",
    ".ts",
    ".css",
    ".sql",
    ".r",
    ".tex",
    ".ini",
    ".cfg",
    ".toml",
    ".sh",
    ".bat",
}

_IMAGE_MIME = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}

_OFFICE_SUFFIXES = {
    ".docx",
    ".docm",
    ".dotx",
    ".dotm",
    ".xlsx",
    ".xlsm",
    ".xltx",
    ".xltm",
    ".pptx",
    ".pptm",
    ".potx",
    ".potm",
    ".odt",
    ".ods",
    ".odp",
}

_DOCX_EXTS = {".docx", ".docm", ".dotx", ".dotm"}
_PPTX_EXTS = {".pptx", ".pptm", ".potx", ".potm"}
_XLSX_EXTS = {".xlsx", ".xlsm", ".xltx", ".xltm"}

_EMBEDDED_IMAGE_MIME = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
}

MAX_TEXT_BYTES = 1_000_000  # ~250K tokens
MAX_IMAGE_BYTES = 5_000_000
MAX_EMBEDDED_IMAGE_BYTES = 400_000
MAX_PDF_PAGES = 30
MAX_PPTX_SLIDES = 30
MAX_XLSX_SHEETS = 10
# Host MCP clients cap tool results at 1MB; keep image bytes well below the cap.
PDF_RESPONSE_BUDGET_BYTES = 650_000
PDF_RENDER_DPI = 110  # lower DPI = smaller output; 110 keeps text readable
PDF_JPEG_QUALITY = 75


def _extract_docx_markdown(filepath: str) -> str:
    """Extract a .docx as GitHub-flavored markdown via pandoc.

    Uses `pypandoc-binary`, which bundles pandoc as a wheel — no system
    install required. Headings, lists, tables, math blocks, and tracked
    changes are preserved; figures become image references, inlined
    separately by `_extract_docx_images`.
    """
    import pypandoc

    return pypandoc.convert_file(
        filepath,
        "gfm",
        format="docx",
        extra_args=["--track-changes=all", "--wrap=none"],
    )


def _extract_docx_images(filepath: str, budget_remaining: int) -> tuple[list, int]:
    """Pull embedded images out of a .docx and return them as ImageContent blocks.

    The .docx zip stores pictures under `word/media/`. We emit them in the
    order they appear there, stopping once `budget_remaining` bytes are used.
    Returns `(blocks, bytes_used)`.
    """
    import base64
    import os as _os
    import zipfile

    from mcp.types import ImageContent

    blocks: list = []
    used = 0
    with zipfile.ZipFile(filepath) as z:
        names = sorted(n for n in z.namelist() if n.startswith("word/media/"))
        for name in names:
            ext = _os.path.splitext(name)[1].lower()
            mime = _EMBEDDED_IMAGE_MIME.get(ext)
            if not mime:
                continue
            blob = z.read(name)
            if len(blob) > MAX_EMBEDDED_IMAGE_BYTES:
                continue
            if used + len(blob) > budget_remaining:
                break
            used += len(blob)
            blocks.append(
                ImageContent(
                    type="image",
                    data=base64.b64encode(blob).decode(),
                    mimeType=mime,
                )
            )
    return blocks, used


def _extract_pptx_blocks(filepath: str, budget: int, pages: str | None = None) -> list:
    """Return interleaved text + image blocks for a .pptx using python-pptx.

    One TextContent per slide (title + body text + speaker notes if present),
    followed by any ImageContent blocks for pictures on that slide. Image
    bytes are tracked against `budget`; once exhausted, later pictures are
    skipped and counted in the trailing summary.

    Args:
        filepath: path to the .pptx file.
        budget: maximum total image bytes to include.
        pages: 1-indexed slide selection like "1-5,7" (parity with PDF
            `pages`). None → first `MAX_PPTX_SLIDES` slides.
    """
    import base64

    from mcp.types import ImageContent, TextContent
    from pptx import Presentation

    prs = Presentation(filepath)
    all_slides = list(prs.slides)
    total = len(all_slides)

    if pages:
        indices = _parse_pages(pages, total)
        truncated_by_cap = False
    else:
        indices = list(range(min(total, MAX_PPTX_SLIDES)))
        truncated_by_cap = total > MAX_PPTX_SLIDES

    blocks: list = []
    used = 0
    skipped_images = 0

    for idx in indices:
        slide = all_slides[idx]
        slide_no = idx + 1
        lines: list[str] = [f"## Slide {slide_no}"]
        images: list[tuple[bytes, str]] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
            try:
                img = shape.image
            except AttributeError, ValueError:
                continue
            images.append((img.blob, img.content_type or "image/png"))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"_Notes: {notes}_")

        blocks.append(TextContent(type="text", text="\n\n".join(lines)))

        for blob, mime in images:
            if len(blob) > MAX_EMBEDDED_IMAGE_BYTES or used + len(blob) > budget:
                skipped_images += 1
                continue
            used += len(blob)
            blocks.append(
                ImageContent(
                    type="image",
                    data=base64.b64encode(blob).decode(),
                    mimeType=mime,
                )
            )

    if truncated_by_cap:
        blocks.append(
            TextContent(
                type="text",
                text=(
                    f"[TRUNCATED — showing first {MAX_PPTX_SLIDES} of "
                    f"{total} slides. Call again with "
                    f"pages='{MAX_PPTX_SLIDES + 1}-{total}' for the rest.]"
                ),
            )
        )
    if skipped_images:
        blocks.append(
            TextContent(
                type="text",
                text=(
                    f"[{skipped_images} image(s) omitted — per-image size cap "
                    f"or response budget reached.]"
                ),
            )
        )
    return blocks


def _extract_xlsx_markdown(
    filepath: str,
    pages: str | None = None,
    max_rows_per_sheet: int = 200,
) -> str:
    """Render sheets of an .xlsx as markdown tables using openpyxl.

    Skips images (spreadsheets rarely use embedded pictures for meaning).
    Truncates each sheet to `max_rows_per_sheet` rows.

    Args:
        filepath: path to the .xlsx file.
        pages: 1-indexed sheet selection like "1-3,5" (parity with PDF
            `pages`). None → first `MAX_XLSX_SHEETS` sheets.
        max_rows_per_sheet: per-sheet row cap.
    """
    from openpyxl import load_workbook

    wb = load_workbook(filepath, data_only=True, read_only=True)
    try:
        all_sheets = wb.worksheets
        total = len(all_sheets)

        if pages:
            indices = _parse_pages(pages, total)
            truncated_by_cap = False
        else:
            indices = list(range(min(total, MAX_XLSX_SHEETS)))
            truncated_by_cap = total > MAX_XLSX_SHEETS

        parts: list[str] = []
        for idx in indices:
            ws = all_sheets[idx]
            rows: list[tuple] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows_per_sheet:
                    break
                rows.append(row)

            max_cols = max((len(r) for r in rows), default=0)

            def col_empty(col: int, rows=rows) -> bool:
                return all(col >= len(r) or r[col] is None for r in rows)

            while max_cols > 0 and col_empty(max_cols - 1):
                max_cols -= 1

            parts.append(f"## Sheet: {ws.title}")
            if max_cols == 0 or not rows:
                parts.append("_(empty)_")
                continue

            def fmt(v: object) -> str:
                if v is None:
                    return ""
                return str(v).replace("|", "\\|").replace("\n", " ")

            def row_line(r: tuple, n: int = max_cols) -> str:
                cells = [fmt(r[i]) if i < len(r) else "" for i in range(n)]
                return "| " + " | ".join(cells) + " |"

            lines = [row_line(rows[0]), "| " + " | ".join(["---"] * max_cols) + " |"]
            for r in rows[1:]:
                lines.append(row_line(r))

            if ws.max_row and ws.max_row > max_rows_per_sheet:
                lines.append(
                    f"\n_[TRUNCATED — showing first {max_rows_per_sheet} of {ws.max_row} rows]_"
                )
            parts.append("\n".join(lines))

        if truncated_by_cap:
            parts.append(
                f"_[TRUNCATED — showing first {MAX_XLSX_SHEETS} of "
                f"{total} sheets. Call again with "
                f"pages='{MAX_XLSX_SHEETS + 1}-{total}' for the rest.]_"
            )
        return "\n\n".join(parts)
    finally:
        wb.close()


def _parse_pages(pages: str | None, total: int) -> list[int]:
    """Parse a page range string like '1-5,7,9-11' into a sorted list of 0-indexed ints.

    Returns all pages if `pages` is None or empty.
    """
    if not pages:
        return list(range(total))
    result: set[int] = set()
    for part in pages.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start, end = part.split("-", 1)
            for p in range(int(start), int(end) + 1):
                if 1 <= p <= total:
                    result.add(p - 1)
        else:
            p = int(part)
            if 1 <= p <= total:
                result.add(p - 1)
    return sorted(result)


def _file_to_content(
    filepath: str,
    pages: str | None = None,
    display_path: str | None = None,
):
    """Convert a file to MCP content blocks Claude.ai can render.

    Args:
        filepath: Path to the file
        pages: For PDFs only — page range like "1-5,7" (1-indexed). None = first MAX_PDF_PAGES.
        display_path: Path shown to the user in headers. Used when the real
            `filepath` is a temp conversion artifact but we want the original
            file's path to appear in the response.

    Returns either a single content block or a list (for PDFs with multiple pages).
    """
    import base64
    import os as _os

    from mcp.types import ImageContent, TextContent

    shown_path = display_path or filepath
    filename = _os.path.basename(shown_path)
    ext = _os.path.splitext(filepath)[1].lower()
    size = _os.path.getsize(filepath)

    if ext in _TEXT_SUFFIXES:
        truncated = False
        with open(filepath, errors="replace") as f:
            text = f.read(MAX_TEXT_BYTES + 1)
        if len(text) > MAX_TEXT_BYTES:
            text = text[:MAX_TEXT_BYTES]
            truncated = True
        header = f"# {filename}\nLocal path: {shown_path}"
        if truncated:
            header += (
                f"\n[TRUNCATED — showing first {MAX_TEXT_BYTES} of {size} bytes. "
                f"Use your host's local file reader on the path above for the full file.]"
            )
        return TextContent(type="text", text=f"{header}\n\n{text}")

    if ext in _IMAGE_MIME:
        if size > MAX_IMAGE_BYTES:
            return TextContent(
                type="text",
                text=f"# {filename}\n[Image too large: {size} bytes > {MAX_IMAGE_BYTES} limit]",
            )
        with open(filepath, "rb") as f:
            data = base64.b64encode(f.read()).decode()
        return ImageContent(type="image", data=data, mimeType=_IMAGE_MIME[ext])

    if ext == ".pdf":
        import pymupdf

        doc = pymupdf.open(filepath)
        total_pages = len(doc)

        if pages:
            page_indices = _parse_pages(pages, total_pages)
            header = (
                f"# {filename} ({total_pages} pages, rendering: {pages})\nLocal path: {shown_path}"
            )
        else:
            page_indices = list(range(min(total_pages, MAX_PDF_PAGES)))
            header = f"# {filename} ({total_pages} pages)\nLocal path: {shown_path}"
            if len(page_indices) < total_pages:
                header += (
                    f"\n[TRUNCATED — rendering first {len(page_indices)} of {total_pages} pages. "
                    f"Call again with pages='X-Y' to get specific pages.]"
                )

        blocks: list = [TextContent(type="text", text=header)]
        # Fit as many requested pages as possible into the response budget.
        # JPEG is much smaller than PNG for document scans / dense text.
        used = 0
        rendered = 0
        skipped_first: int | None = None
        for i in page_indices:
            pix = doc[i].get_pixmap(dpi=PDF_RENDER_DPI)
            img_bytes = pix.tobytes("jpeg", jpg_quality=PDF_JPEG_QUALITY)
            if used + len(img_bytes) > PDF_RESPONSE_BUDGET_BYTES and rendered > 0:
                skipped_first = i + 1  # 1-indexed page number
                break
            used += len(img_bytes)
            rendered += 1
            blocks.append(
                ImageContent(
                    type="image",
                    data=base64.b64encode(img_bytes).decode(),
                    mimeType="image/jpeg",
                )
            )
        doc.close()
        if skipped_first is not None:
            last_rendered = page_indices[rendered - 1] + 1
            blocks.append(
                TextContent(
                    type="text",
                    text=(
                        f"[TRUNCATED — fit {rendered} pages into the response "
                        f"size budget. Rendered through page {last_rendered}. "
                        f"Call again with pages='{skipped_first}-...' to get "
                        f"the next pages.]"
                    ),
                )
            )
        return blocks

    if ext in _DOCX_EXTS:
        md = _extract_docx_markdown(filepath)
        truncated = len(md) > MAX_TEXT_BYTES
        if truncated:
            md = md[:MAX_TEXT_BYTES]
        header = f"# {filename}\nLocal path: {shown_path}"
        if truncated:
            header += f"\n[TRUNCATED — showing first {MAX_TEXT_BYTES} chars of extracted markdown.]"
        text_block = TextContent(type="text", text=f"{header}\n\n{md}")
        image_budget = max(0, PDF_RESPONSE_BUDGET_BYTES - len(text_block.text))
        blocks: list = [text_block]
        image_blocks, _ = _extract_docx_images(
            filepath,
            budget_remaining=image_budget,
        )
        blocks.extend(image_blocks)
        return blocks

    if ext in _PPTX_EXTS:
        header = f"# {filename}\nLocal path: {shown_path}"
        header_block = TextContent(type="text", text=header)
        image_budget = max(0, PDF_RESPONSE_BUDGET_BYTES - len(header))
        blocks = [header_block]
        blocks.extend(_extract_pptx_blocks(filepath, budget=image_budget, pages=pages))
        return blocks

    if ext in _XLSX_EXTS:
        md = _extract_xlsx_markdown(filepath, pages=pages)
        truncated = len(md) > MAX_TEXT_BYTES
        if truncated:
            md = md[:MAX_TEXT_BYTES]
        header = f"# {filename}\nLocal path: {shown_path}"
        if truncated:
            header += f"\n[TRUNCATED — showing first {MAX_TEXT_BYTES} chars of extracted markdown.]"
        return TextContent(type="text", text=f"{header}\n\n{md}")

    return TextContent(
        type="text",
        text=(
            f"# {filename}\n"
            f"Size: {size} bytes\n"
            f"Absolute path on the user's local machine: {shown_path}\n"
            f"\n"
            f"This is a SINGLE BINARY FILE saved at the absolute path above;\n"
            f"the server did not extract it. Do NOT call `read_downloaded_file`\n"
            f"on it — that tool is only for paths returned by a previous\n"
            f"zip-listing response.\n"
            f"\n"
            f"If your host agent has a local file-read tool (Claude Code,\n"
            f"Cursor, Cline, Continue, Aider and similar all ship one), call\n"
            f"it on the absolute path above. Otherwise, ask the user to\n"
            f"upload or paste the relevant content."
        ),
    )


@mcp.tool()
async def download_resource(url: str, pages: str | None = None):
    """Download a Moodle file and return its content.

    Accepts either a direct pluginfile URL OR a `/mod/resource/view.php?id=...`
    URL — redirects are followed automatically. For any module of type
    "resource" from `get_course_contents`, call this tool directly with the
    module's URL. Do NOT call `get_module_content` on resource URLs first;
    those pages 303-redirect straight to the file and will error.

    - Text files: returned as text (max 1MB, truncated if larger)
    - Images (png/jpg/gif/webp): returned inline (max 5MB)
    - PDFs: rasterized to JPEG images. Default = first 30 pages, but fewer
      may be returned to fit the host tool-result size cap — the response
      will tell you the next page to request if truncated.
      Use `pages` to select specific pages (e.g. "1-5,7,10-12").
    - Zips: extracted but NOT auto-loaded — returns a file listing instead.
      Use `read_downloaded_file` to load specific files from the zip.
    - Word docs (.docx/.docm): text converted to GitHub-flavored markdown
      via pandoc (headings, lists, tables, tracked changes, math blocks),
      plus embedded images from `word/media/` inlined as image blocks.
      `pages` is ignored for .docx (no natural pagination).
    - PowerPoint (.pptx/.pptm): one text block per slide (title, body,
      speaker notes) followed by that slide's pictures as image blocks.
      Default = first 30 slides; use `pages` to pick specific slides
      (e.g. "1-10,15").
    - Excel (.xlsx/.xlsm): each selected sheet as a markdown table (first
      200 rows per sheet). Default = first 10 sheets; use `pages` to pick
      specific sheets by 1-indexed position (e.g. "1,3-4").
    - ODF formats (.odt/.ods/.odp) and other binaries: returned as a
      pointer notice with the file's absolute path.

    Args:
        url: Moodle resource or pluginfile URL
        pages: 1-indexed selection like "1-5,7,10-12". Meaning depends on
            format: PDF pages, .pptx slides, or .xlsx sheet positions.
            Ignored for .docx and non-paginated formats.
    """
    import os as _os
    import zipfile

    from mcp.types import TextContent

    filepath = await download_file(url)
    filename = _os.path.basename(filepath)

    # this is because office formats are technically zip containers
    ext = _os.path.splitext(filename)[1].lower()
    is_office = ext in _OFFICE_SUFFIXES

    if not is_office and zipfile.is_zipfile(filepath):
        extract_dir = _os.path.join(_os.path.dirname(filepath), _os.path.splitext(filename)[0])
        _os.makedirs(extract_dir, exist_ok=True)
        with zipfile.ZipFile(filepath) as zf:
            zf.extractall(extract_dir)
        _os.remove(filepath)

        listing_lines = [f"# {filename} (zip archive — choose files to load)"]
        listing_lines.append("Use `read_downloaded_file(path=...)` with one of these paths:")
        listing_lines.append("")
        for root, _, names in _os.walk(extract_dir):
            for name in sorted(names):
                fpath = _os.path.join(root, name)
                rel = _os.path.relpath(fpath, _os.path.dirname(filepath))
                size = _os.path.getsize(fpath)
                listing_lines.append(f"- `{rel}` ({size} bytes)")

        return [TextContent(type="text", text="\n".join(listing_lines))]

    content = _file_to_content(filepath, pages=pages)
    if isinstance(content, list):
        return content
    return [content]


@mcp.tool()
async def read_downloaded_file(path: str, pages: str | None = None):
    """Read a previously downloaded file from the local downloads directory.

    ONLY use this after `download_resource` returned a zip-listing response
    (a "zip archive — choose files to load" header with bullet-listed paths).
    Pass one of those exact listed paths verbatim — do NOT invent, modify,
    or guess paths.

    Do NOT use this tool for any path not literally present in a prior zip
    listing. Office documents are not extracted as zip trees — call
    `download_resource` on their URL instead; it handles .docx/.pptx/.xlsx
    content directly.

    Args:
        path: Relative path under the downloads dir, copied verbatim from a
            zip listing returned by `download_resource`.
        pages: 1-indexed selection like "1-5,7" — PDF pages, .pptx slides,
            or .xlsx sheet positions. Ignored for other formats.
    """
    import os as _os

    from moodler_mcp.client import DOWNLOADS_DIR

    # Security: prevent path traversal
    if ".." in path or path.startswith("/"):
        raise ValueError("Invalid path")

    filepath = _os.path.join(DOWNLOADS_DIR, path)
    if not _os.path.exists(filepath):
        raise FileNotFoundError(f"Not found: {path}")

    content = _file_to_content(filepath, pages=pages)
    if isinstance(content, list):
        return content
    return [content]


@mcp.tool()
async def get_module_content(url: str) -> str:
    """Get content from a Moodle module page (assignment, folder, URL, page, etc).

    Extracts files, links, and descriptions from any Moodle module page.
    Use download_resource to download individual files from the results.

    IMPORTANT: Do NOT call this on `/mod/resource/view.php?id=...` URLs
    (modules of type "resource" from `get_course_contents`). Those pages
    303-redirect directly to the underlying file and have no HTML content
    to scrape — call `download_resource(url)` on the resource URL instead.
    This tool is for `/mod/assign/`, `/mod/folder/`, `/mod/url/`, `/mod/page/`,
    and similar wrapper pages.

    Args:
        url: Moodle module URL (e.g. '/mod/assign/view.php?id=570007')
    """
    from urllib.parse import urlparse

    from bs4 import BeautifulSoup

    if not url.startswith("http"):
        from moodler_mcp.config import MOODLE_URL as base

        url = f"{base}{url}"

    parsed = urlparse(url)
    path = parsed.path
    html = await fetch_page(f"{parsed.path}?{parsed.query}" if parsed.query else parsed.path)
    soup = BeautifulSoup(html, "html.parser")

    result: dict = {}

    title_el = soup.select_one("h2, .page-header-headings h1")
    if title_el:
        result["title"] = title_el.get_text(strip=True)

    intro = soup.select_one(".intro, .activity-description, .assignmentintro, .mod-description")
    if intro:
        result["description"] = intro.get_text(strip=True)

    files = []
    for a in soup.select("a[href*='pluginfile']"):
        name = a.get_text(strip=True)
        href = a.get("href", "")
        if name and href:
            files.append({"name": name, "url": href})
    if files:
        result["files"] = files

    if "/mod/url/" in path:
        for a in soup.select(".urlworkaround a, .resourceworkaround a"):
            result["external_url"] = a.get("href", "")
            break

    if "/mod/folder/" in path:
        folder_files = []
        for a in soup.select(".fp-filename-icon a, .foldertree a[href*='pluginfile']"):
            name = a.get_text(strip=True)
            href = a.get("href", "")
            if name and href:
                folder_files.append({"name": name, "url": href})
        if folder_files:
            result["files"] = folder_files

    if "/mod/assign/" in path:
        for row in soup.select(".submissionstatustable tr, .generaltable tr"):
            cells = row.select("td")
            if len(cells) >= 2:
                label = cells[0].get_text(strip=True).lower()
                value = cells[1].get_text(strip=True)
                if "due" in label:
                    result["due_date"] = value
                elif "status" in label and "submission" in label:
                    result["submission_status"] = value
                elif "grading" in label and "status" in label:
                    result["grading_status"] = value
                elif "grade" in label and "status" not in label:
                    result["grade"] = value

    return json.dumps(result, indent=2)
